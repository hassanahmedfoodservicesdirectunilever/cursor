from datetime import datetime
from math import atan2, cos, sin
from pathlib import Path
from shutil import copyfile
from textwrap import fill

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).parent
ASSETS_DIR = ROOT / "presentation_assets"
LEGACY_OUTPUT_FILE = ROOT / "Solution_Assessment_Cursor_AI_Capabilities_Integration_Review.pptx"
PARTICIPANT_OUTPUT_FILE = ROOT / "Solution_Assessment_Cursor_AI_Capabilities_Integration_Review_Participant.pptx"
TRAINER_OUTPUT_FILE = ROOT / "Solution_Assessment_Cursor_AI_Capabilities_Integration_Review_Trainer_45min.pptx"

CANVAS_W = 2400
CANVAS_H = 1350


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def save_image(image: Image.Image, path: Path) -> None:
    image.save(path, format="PNG", dpi=(300, 300))


def draw_header(draw: ImageDraw.ImageDraw, title: str, width: int, accent=(26, 92, 168)) -> None:
    draw.rectangle((0, 0, width, 140), fill=accent)
    draw.text((48, 38), title, font=load_font(64, bold=True), fill=(255, 255, 255))


def new_canvas(title: str, accent=(26, 92, 168)):
    image = Image.new("RGB", (CANVAS_W, CANVAS_H), (244, 248, 255))
    draw = ImageDraw.Draw(image)
    draw_header(draw, title, CANVAS_W, accent=accent)
    return image, draw


def draw_arrow(draw: ImageDraw.ImageDraw, start, end, color=(70, 109, 175), width=8) -> None:
    draw.line([start, end], fill=color, width=width)
    angle = atan2(end[1] - start[1], end[0] - start[0])
    arrow_len = 28
    arrow_width = 14
    p1 = end
    p2 = (
        end[0] - arrow_len * cos(angle) + arrow_width * sin(angle),
        end[1] - arrow_len * sin(angle) - arrow_width * cos(angle),
    )
    p3 = (
        end[0] - arrow_len * cos(angle) - arrow_width * sin(angle),
        end[1] - arrow_len * sin(angle) + arrow_width * cos(angle),
    )
    draw.polygon([p1, p2, p3], fill=color)


def draw_text_center(draw: ImageDraw.ImageDraw, box, text: str, font, fill_color=(24, 45, 78)) -> None:
    x1, y1, x2, y2 = box
    bb = draw.textbbox((0, 0), text, font=font)
    text_w = bb[2] - bb[0]
    text_h = bb[3] - bb[1]
    draw.text((x1 + (x2 - x1 - text_w) / 2, y1 + (y2 - y1 - text_h) / 2), text, font=font, fill=fill_color)


def draw_card(draw: ImageDraw.ImageDraw, box, title: str, lines, fill_color=(236, 244, 255), border=(130, 161, 208)) -> None:
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=24, fill=fill_color, outline=border, width=4)
    draw.text((x1 + 24, y1 + 20), title, font=load_font(44, bold=True), fill=(23, 45, 80))
    y = y1 + 88
    for line in lines:
        wrapped = fill(line, width=29)
        draw.text((x1 + 26, y), f"- {wrapped}", font=load_font(30), fill=(39, 58, 92))
        y += 62 + (wrapped.count("\n") * 18)


def create_cover_visual(path: Path) -> None:
    image, draw = new_canvas("Cursor AI + MCP: simple view")

    center = (1200, 700)
    r = 250
    draw.ellipse((center[0] - r, center[1] - r, center[0] + r, center[1] + r), fill=(28, 83, 160), outline=(17, 58, 120), width=6)
    draw_text_center(draw, (center[0] - r, center[1] - r, center[0] + r, center[1] + r), "Cursor\nAgent", load_font(64, bold=True), (255, 255, 255))

    nodes = [
        ((320, 300, 760, 500), "Jira MCP", (255, 245, 230)),
        ((1640, 300, 2080, 500), "Figma MCP", (235, 250, 241)),
        ((320, 900, 760, 1100), "Bitbucket MCP", (236, 244, 255)),
        ((1640, 900, 2080, 1100), "Team Rules", (251, 240, 250)),
    ]
    for box, label, color in nodes:
        draw.rounded_rectangle(box, radius=26, fill=color, outline=(136, 166, 211), width=4)
        draw_text_center(draw, box, label, load_font(52, bold=True))

    draw_arrow(draw, (760, 400), (940, 560))
    draw_arrow(draw, (1640, 400), (1460, 560))
    draw_arrow(draw, (760, 1000), (940, 840))
    draw_arrow(draw, (1640, 1000), (1460, 840))

    draw.text((48, 1220), "Simple goal: connect tools, use safe prompts, and speed up delivery.", font=load_font(48, bold=True), fill=(39, 60, 94))
    save_image(image, path)


def create_top_mcp_visual(path: Path) -> None:
    image, draw = new_canvas("Top MCP tools for development teams", accent=(18, 99, 151))
    draw.text((58, 170), "Start with these first. They give fast value in engineering work.", font=load_font(42), fill=(40, 60, 95))

    cards = [
        ("1) Bitbucket / Git MCP", ["Read repos and pull requests", "Draft PR summaries", "Build better review checklists"]),
        ("2) Jira MCP", ["Read issues and sprint data", "Write clear acceptance criteria", "Create standup summaries"]),
        ("3) Figma MCP", ["Read design tokens", "Map components to stories", "Reduce design-code mismatch"]),
        ("4) Docs MCP", ["Search team docs quickly", "Use standards in prompts", "Reduce repeated questions"]),
        ("5) CI/CD MCP", ["Read build status quickly", "Summarize failures", "Suggest fix order"]),
        ("6) Database MCP", ["Read schema and queries", "Check impact before code changes", "Improve migration safety"]),
    ]

    x = 72
    y = 260
    w = 740
    h = 430
    gap_x = 54
    gap_y = 44
    for idx, (title, lines) in enumerate(cards):
        col = idx % 3
        row = idx // 3
        box = (x + col * (w + gap_x), y + row * (h + gap_y), x + col * (w + gap_x) + w, y + row * (h + gap_y) + h)
        fill_color = (237, 245, 255) if row == 0 else (238, 251, 243)
        draw_card(draw, box, title, lines, fill_color=fill_color)

    save_image(image, path)


def create_top_skills_visual(path: Path) -> None:
    image, draw = new_canvas("Top agent skills to start with", accent=(83, 90, 180))
    draw.text((58, 170), "Build these 6 reusable skills first for quick adoption.", font=load_font(42), fill=(44, 58, 100))

    skills = [
        ("jira-ticket-triage", "Input: issue key\nOutput: summary + test cases"),
        ("figma-handoff", "Input: frame id\nOutput: dev task checklist"),
        ("pr-quality-check", "Input: PR id\nOutput: review checklist"),
        ("release-note-writer", "Input: commit range\nOutput: release notes"),
        ("bug-root-cause", "Input: bug + logs\nOutput: likely causes + next steps"),
        ("test-case-generator", "Input: user story\nOutput: unit + API test ideas"),
    ]

    x = 76
    y = 250
    w = 730
    h = 430
    gap_x = 54
    gap_y = 44
    for idx, (name, desc) in enumerate(skills):
        col = idx % 3
        row = idx // 3
        box = (x + col * (w + gap_x), y + row * (h + gap_y), x + col * (w + gap_x) + w, y + row * (h + gap_y) + h)
        fill_color = (244, 240, 255) if row == 0 else (236, 247, 255)
        draw.rounded_rectangle(box, radius=24, fill=fill_color, outline=(146, 164, 210), width=4)
        draw.text((box[0] + 24, box[1] + 26), name, font=load_font(42, bold=True), fill=(33, 53, 96))
        draw.text((box[0] + 24, box[1] + 120), desc, font=load_font(34), fill=(45, 65, 99))
        draw.rectangle((box[0] + 24, box[1] + 290, box[2] - 24, box[1] + 362), fill=(255, 247, 230), outline=(214, 166, 104), width=3)
        draw.text((box[0] + 38, box[1] + 312), "Save as versioned skill package", font=load_font(31, bold=True), fill=(105, 68, 30))

    save_image(image, path)


def create_architecture_visual(path: Path) -> None:
    image, draw = new_canvas("Simple architecture for safe use", accent=(26, 116, 142))

    draw_card(
        draw,
        (120, 260, 620, 1050),
        "Tools",
        ["Jira", "Figma", "Bitbucket", "Docs and CI"],
        fill_color=(238, 247, 255),
    )
    draw_card(
        draw,
        (760, 260, 1260, 1050),
        "MCP layer",
        ["API auth", "Tool schema", "Logs and rate limits", "Error handling"],
        fill_color=(236, 252, 243),
    )
    draw_card(
        draw,
        (1400, 260, 1900, 1050),
        "Cursor",
        ["Prompt + skill", "User approval", "Code and docs output", "Action summary"],
        fill_color=(245, 241, 255),
    )
    draw_card(
        draw,
        (1960, 260, 2320, 1050),
        "Governance",
        ["Least privilege", "Audit logs", "Weekly KPI check", "Rotate tokens"],
        fill_color=(255, 245, 236),
    )

    draw_arrow(draw, (620, 620), (760, 620), width=10)
    draw_arrow(draw, (1260, 620), (1400, 620), width=10)
    draw_arrow(draw, (1900, 620), (1960, 620), width=10)

    save_image(image, path)


def draw_fake_cursor_shell(draw: ImageDraw.ImageDraw, title: str) -> None:
    draw.rounded_rectangle((220, 190, 2180, 1200), radius=24, fill=(23, 28, 38), outline=(75, 87, 112), width=4)
    draw.rectangle((220, 190, 2180, 260), fill=(34, 41, 57))
    draw.text((270, 210), "Cursor - Example UI", font=load_font(34, bold=True), fill=(226, 233, 247))
    draw.text((790, 210), title, font=load_font(34), fill=(170, 184, 212))

    for i, c in enumerate([(250, 106, 106), (245, 183, 78), (93, 208, 105)]):
        draw.ellipse((236 + i * 30, 213, 256 + i * 30, 233), fill=c)


def create_cursor_settings_screen(path: Path) -> None:
    image, draw = new_canvas("Screenshot: open MCP settings inside Cursor", accent=(40, 100, 158))
    draw_fake_cursor_shell(draw, "Settings")

    draw.rectangle((260, 280, 620, 1160), fill=(29, 35, 48))
    menu_items = ["General", "Editor", "Terminal", "Features", "MCP", "Privacy", "About"]
    y = 330
    for item in menu_items:
        if item == "MCP":
            draw.rounded_rectangle((290, y - 14, 590, y + 40), radius=12, fill=(64, 104, 184))
            draw.text((320, y), item, font=load_font(36, bold=True), fill=(255, 255, 255))
        else:
            draw.text((320, y), item, font=load_font(34), fill=(195, 207, 231))
        y += 92

    draw.rectangle((650, 280, 2140, 1160), fill=(37, 45, 61))
    draw.text((710, 340), "MCP Servers", font=load_font(46, bold=True), fill=(233, 239, 252))
    draw.rounded_rectangle((700, 430, 2060, 540), radius=14, fill=(50, 62, 84))
    draw.text((740, 466), "Enable MCP integration", font=load_font(36), fill=(224, 231, 246))
    draw.rounded_rectangle((1840, 452, 2020, 520), radius=34, fill=(68, 158, 91))
    draw.text((1886, 468), "ON", font=load_font(34, bold=True), fill=(255, 255, 255))

    draw.rounded_rectangle((700, 590, 2060, 730), radius=14, fill=(50, 62, 84))
    draw.text((740, 635), "Open mcp.json", font=load_font(36), fill=(224, 231, 246))
    draw.rounded_rectangle((1710, 614, 2020, 694), radius=12, fill=(73, 122, 220))
    draw.text((1766, 635), "Open file", font=load_font(34, bold=True), fill=(255, 255, 255))

    draw.rounded_rectangle((70, 1020, 1020, 1270), radius=18, fill=(237, 246, 255), outline=(130, 161, 207), width=3)
    draw.text((108, 1062), "Step 1: Click Settings  ->  Features  ->  MCP", font=load_font(40, bold=True), fill=(30, 58, 101))
    draw.text((108, 1130), "Step 2: Turn ON MCP and open mcp.json", font=load_font(36), fill=(30, 58, 101))

    save_image(image, path)


def create_cursor_mcp_json_screen(path: Path) -> None:
    image, draw = new_canvas("Screenshot: mcp.json setup", accent=(35, 119, 170))
    draw_fake_cursor_shell(draw, "mcp.json")

    draw.rectangle((260, 280, 650, 1160), fill=(28, 36, 49))
    draw.text((300, 330), "Explorer", font=load_font(34, bold=True), fill=(220, 229, 246))
    files = [".cursor/", "mcp.json", ".env", "skills/", "README.md"]
    y = 405
    for item in files:
        col = (255, 255, 255) if item == "mcp.json" else (192, 205, 230)
        draw.text((320, y), item, font=load_font(33, bold=(item == "mcp.json")), fill=col)
        y += 70

    draw.rectangle((680, 280, 2140, 1160), fill=(18, 24, 35))
    json_lines = [
        "{",
        '  "mcpServers": {',
        '    "jira": {',
        '      "command": "python",',
        '      "args": ["servers/jira_server.py"],',
        '      "envFile": ".env"',
        "    },",
        '    "figma": {',
        '      "command": "python",',
        '      "args": ["servers/figma_server.py"],',
        '      "envFile": ".env"',
        "    },",
        '    "bitbucket": {',
        '      "command": "python",',
        '      "args": ["servers/bitbucket_server.py"],',
        '      "envFile": ".env"',
        "    }",
        "  }",
        "}",
    ]
    y = 330
    for idx, line in enumerate(json_lines, start=1):
        draw.text((720, y), f"{idx:>2}", font=load_font(30), fill=(107, 122, 149))
        draw.text((790, y), line, font=load_font(32), fill=(220, 231, 248))
        y += 45

    draw.rounded_rectangle((1760, 1040, 2060, 1130), radius=12, fill=(72, 126, 218))
    draw.text((1835, 1068), "Save", font=load_font(36, bold=True), fill=(255, 255, 255))

    draw.rounded_rectangle((70, 1020, 1130, 1270), radius=18, fill=(238, 252, 243), outline=(120, 176, 145), width=3)
    draw.text((108, 1062), "Step 3: Paste config and save file.", font=load_font(40, bold=True), fill=(28, 91, 60))
    draw.text((108, 1130), "Step 4: Restart Cursor to load servers.", font=load_font(36), fill=(28, 91, 60))

    save_image(image, path)


def create_cursor_connection_status_screen(path: Path) -> None:
    image, draw = new_canvas("Screenshot: verify connections", accent=(24, 135, 96))
    draw_fake_cursor_shell(draw, "MCP connection status")

    draw.rectangle((260, 300, 2140, 1150), fill=(34, 44, 61))
    draw.text((320, 360), "MCP Server Status", font=load_font(46, bold=True), fill=(235, 241, 253))

    rows = [
        ("jira", "Connected", "Read issues, create summaries"),
        ("figma", "Connected", "Read design tokens and frames"),
        ("bitbucket", "Connected", "Read repos and pull requests"),
    ]
    y = 470
    for name, status, detail in rows:
        draw.rounded_rectangle((320, y, 2080, y + 170), radius=16, fill=(48, 61, 84))
        draw.text((380, y + 45), name, font=load_font(40, bold=True), fill=(230, 238, 252))
        draw.rounded_rectangle((760, y + 42, 1020, y + 118), radius=12, fill=(66, 162, 97))
        draw.text((806, y + 62), status, font=load_font(34, bold=True), fill=(255, 255, 255))
        draw.text((1080, y + 58), detail, font=load_font(34), fill=(214, 225, 246))
        draw.rounded_rectangle((1850, y + 46, 2030, y + 120), radius=12, fill=(76, 129, 225))
        draw.text((1890, y + 64), "Test", font=load_font(34, bold=True), fill=(255, 255, 255))
        y += 210

    draw.rounded_rectangle((70, 1010, 1300, 1270), radius=18, fill=(234, 247, 255), outline=(128, 162, 209), width=3)
    draw.text((110, 1052), "Step 5: Click Test on each server.", font=load_font(40, bold=True), fill=(35, 59, 99))
    draw.text((110, 1120), "Step 6: Use a simple prompt to check output.", font=load_font(36), fill=(35, 59, 99))

    save_image(image, path)


def create_tutorial_path_visual(path: Path) -> None:
    image, draw = new_canvas("Spoon-feed tutorial path", accent=(59, 102, 182))
    steps = [
        "1) Install tools",
        "2) Add tokens",
        "3) Open MCP settings",
        "4) Add mcp.json",
        "5) Test Jira",
        "6) Test Figma",
        "7) Test Bitbucket",
        "8) Build first skill",
    ]
    x = 120
    y = 360
    w = 250
    h = 240
    gap = 30
    for idx, step in enumerate(steps):
        box = (x, y, x + w, y + h)
        draw.rounded_rectangle(box, radius=18, fill=(237, 245, 255), outline=(126, 157, 208), width=3)
        draw.text((x + 18, y + 55), fill(step, width=16), font=load_font(33, bold=True), fill=(27, 51, 91))
        if idx < len(steps) - 1:
            draw_arrow(draw, (x + w, y + h // 2), (x + w + gap - 6, y + h // 2))
        x += w + gap

    draw.rounded_rectangle((170, 760, 2230, 1040), radius=20, fill=(236, 252, 243), outline=(126, 180, 149), width=3)
    draw.text((230, 820), "Easy method: trainer demo (10 min) -> pair lab (25 min) -> review and fix (15 min).", font=load_font(42, bold=True), fill=(28, 93, 63))
    draw.text((230, 892), "Everyone should complete one full workflow in the same day.", font=load_font(38), fill=(28, 93, 63))

    save_image(image, path)


def create_daily_workflow_visual(path: Path) -> None:
    image, draw = new_canvas("Daily team workflow (simple)", accent=(16, 120, 141))
    blocks = [
        ("Morning", ["Open Jira tasks", "Plan with Cursor prompt"]),
        ("Build", ["Write code with Cursor", "Use Figma + Bitbucket context"]),
        ("Review", ["Run PR checklist skill", "Ask agent for risk summary"]),
        ("Close", ["Update Jira status", "Save learning to skill library"]),
    ]
    x = 220
    y = 340
    w = 440
    h = 560
    gap = 120
    colors = [(236, 246, 255), (236, 252, 243), (244, 240, 255), (255, 245, 236)]
    for idx, (title, lines) in enumerate(blocks):
        box = (x, y, x + w, y + h)
        draw_card(draw, box, title, lines, fill_color=colors[idx])
        if idx < len(blocks) - 1:
            draw_arrow(draw, (x + w, y + h // 2), (x + w + gap - 20, y + h // 2), width=10)
        x += w + gap

    save_image(image, path)


def create_risk_controls_visual(path: Path) -> None:
    image, draw = new_canvas("Risk controls in plain English", accent=(42, 98, 172))
    controls = [
        ("Use low access tokens", "Only give needed access.\nDo not use admin tokens."),
        ("Ask approval before write", "Any status change or merge\nneeds human approval."),
        ("Keep logs", "Save prompt, tool call,\nand final output."),
        ("Check weekly KPIs", "Review failures and remove\nunused skills."),
    ]
    x = 120
    y = 320
    w = 520
    h = 420
    gap = 58
    for idx, (title, desc) in enumerate(controls):
        row = idx // 2
        col = idx % 2
        x1 = x + col * (w + gap)
        y1 = y + row * (h + 100)
        box = (x1, y1, x1 + w, y1 + h)
        draw.rounded_rectangle(box, radius=24, fill=(237, 246, 255), outline=(126, 159, 208), width=4)
        draw.ellipse((x1 + 26, y1 + 26, x1 + 120, y1 + 120), fill=(56, 114, 200))
        draw.text((x1 + 61, y1 + 50), str(idx + 1), font=load_font(44, bold=True), fill=(255, 255, 255))
        draw.text((x1 + 140, y1 + 42), title, font=load_font(40, bold=True), fill=(27, 49, 88))
        draw.text((x1 + 42, y1 + 152), desc, font=load_font(34), fill=(42, 62, 96))

    draw.rounded_rectangle((1220, 320, 2270, 840), radius=24, fill=(236, 252, 243), outline=(126, 181, 149), width=4)
    draw.text((1270, 372), "Simple weekly checklist", font=load_font(46, bold=True), fill=(28, 95, 61))
    checklist = [
        "1. Review failed MCP calls.",
        "2. Fix or disable weak prompts.",
        "3. Rotate old tokens.",
        "4. Publish one improved skill.",
    ]
    y = 470
    for item in checklist:
        draw.text((1280, y), item, font=load_font(34), fill=(28, 95, 61))
        y += 86

    save_image(image, path)


def create_roadmap_visual(path: Path) -> None:
    image, draw = new_canvas("30-60-90 day rollout plan", accent=(25, 109, 168))
    draw.line((220, 720, 2180, 720), fill=(73, 112, 177), width=12)
    phases = [
        (500, "Day 0-30", ["Connect Jira/Figma/Bitbucket", "Train first squad", "Set KPI baseline"]),
        (1200, "Day 31-60", ["Publish 6 core skills", "Run weekly office hours", "Track adoption"]),
        (1900, "Day 61-90", ["Scale to more squads", "Improve weak prompts", "Share business impact"]),
    ]
    for x, title, lines in phases:
        draw.ellipse((x - 36, 684, x + 36, 756), fill=(35, 94, 177))
        draw.rounded_rectangle((x - 300, 280, x + 300, 620), radius=22, fill=(235, 244, 255), outline=(128, 159, 207), width=4)
        draw.text((x - 250, 320), title, font=load_font(48, bold=True), fill=(26, 48, 83))
        yy = 410
        for line in lines:
            draw.text((x - 250, yy), f"- {line}", font=load_font(33), fill=(40, 60, 94))
            yy += 70

    draw.rounded_rectangle((260, 860, 2140, 1130), radius=20, fill=(236, 252, 243), outline=(126, 180, 149), width=3)
    draw.text((320, 930), "Success target: 20% faster delivery with safe controls and clear audit logs.", font=load_font(44, bold=True), fill=(26, 92, 60))
    draw.text((320, 995), "Main KPIs: cycle time, PR lead time, reopen rate, prompt reuse.", font=load_font(38), fill=(26, 92, 60))

    save_image(image, path)


def create_prompt_formula_visual(path: Path) -> None:
    image, draw = new_canvas("Prompt formula (use this every time)", accent=(73, 95, 188))

    blocks = [
        ("1) Context", ["What system?", "What ticket?", "What code area?"]),
        ("2) Task", ["What do you want?", "Summary?", "Code?", "Checklist?"]),
        ("3) Constraints", ["No write action", "Use style guide", "Keep output short"]),
        ("4) Output format", ["Bullets", "Table", "JSON", "Commit message"]),
    ]
    x = 180
    y = 340
    w = 480
    h = 640
    gap = 80
    colors = [(236, 246, 255), (236, 252, 243), (244, 240, 255), (255, 245, 236)]
    for idx, (title, lines) in enumerate(blocks):
        box = (x, y, x + w, y + h)
        draw_card(draw, box, title, lines, fill_color=colors[idx])
        if idx < len(blocks) - 1:
            draw_arrow(draw, (x + w, y + h // 2), (x + w + gap - 20, y + h // 2), width=10)
        x += w + gap

    draw.rounded_rectangle((220, 1030, 2180, 1240), radius=20, fill=(236, 252, 243), outline=(124, 180, 148), width=4)
    draw.text(
        (280, 1092),
        "Template: Context + Task + Constraints + Output format = better and repeatable results.",
        font=load_font(40, bold=True),
        fill=(28, 94, 62),
    )

    save_image(image, path)


def create_do_dont_visual(path: Path) -> None:
    image, draw = new_canvas("Do and Don't for safe adoption", accent=(27, 112, 162))

    left_box = (160, 260, 1130, 1130)
    right_box = (1270, 260, 2240, 1130)
    draw.rounded_rectangle(left_box, radius=24, fill=(236, 252, 243), outline=(121, 178, 145), width=4)
    draw.rounded_rectangle(right_box, radius=24, fill=(255, 240, 240), outline=(206, 141, 141), width=4)
    draw.text((220, 320), "DO", font=load_font(64, bold=True), fill=(29, 96, 61))
    draw.text((1330, 320), "DON'T", font=load_font(64, bold=True), fill=(145, 64, 64))

    do_lines = [
        "Use clear and short prompts.",
        "Use read-only mode first.",
        "Ask approval before write actions.",
        "Save good prompts as skills.",
        "Check logs and KPIs every week.",
    ]
    dont_lines = [
        "Do not use admin tokens.",
        "Do not auto-merge without review.",
        "Do not skip test/checklist step.",
        "Do not keep stale skills forever.",
        "Do not hide failed automations.",
    ]

    y = 450
    for line in do_lines:
        draw.text((220, y), f"- {line}", font=load_font(40), fill=(30, 94, 61))
        y += 125

    y = 450
    for line in dont_lines:
        draw.text((1330, y), f"- {line}", font=load_font(40), fill=(145, 64, 64))
        y += 125

    save_image(image, path)


def create_common_errors_visual(path: Path) -> None:
    image, draw = new_canvas("Common errors and quick fixes", accent=(45, 98, 172))
    draw.text((58, 178), "Use this slide during live training when someone gets blocked.", font=load_font(40), fill=(39, 60, 95))

    columns = [150, 780, 1380, 2230]
    top = 260
    bottom = 1130
    draw.rounded_rectangle((columns[0], top, columns[-1], bottom), radius=18, fill=(237, 246, 255), outline=(128, 160, 208), width=4)

    headers = ["Error", "What it means", "Quick fix"]
    for i in range(3):
        draw.rectangle((columns[i], top, columns[i + 1], top + 100), fill=(77, 118, 187))
        draw.text((columns[i] + 24, top + 28), headers[i], font=load_font(40, bold=True), fill=(255, 255, 255))

    rows = [
        ("401 Unauthorized", "Token is wrong or expired", "Create new token and update .env"),
        ("403 Forbidden", "No access to project or file", "Ask for project/file permission"),
        ("404 Not Found", "Wrong URL, key, or workspace", "Check project key, file key, repo name"),
        ("Empty API result", "Filter is too strict", "Try wider query or remove filter"),
        ("MCP not listed", "mcp.json not loaded", "Save file and restart Cursor"),
    ]

    y = top + 130
    for err, meaning, fix in rows:
        row_bottom = y + 145
        draw.line((columns[0], row_bottom, columns[-1], row_bottom), fill=(154, 177, 216), width=2)
        draw.text((columns[0] + 18, y + 30), err, font=load_font(34, bold=True), fill=(34, 54, 89))
        draw.text((columns[1] + 18, y + 30), meaning, font=load_font(34), fill=(34, 54, 89))
        draw.text((columns[2] + 18, y + 30), fix, font=load_font(34), fill=(34, 54, 89))
        y += 145

    save_image(image, path)


def create_five_min_routine_visual(path: Path) -> None:
    image, draw = new_canvas("5-minute daily routine (very easy)", accent=(20, 128, 145))
    steps = [
        ("Minute 1", "Open Jira list and pick top task"),
        ("Minute 2", "Ask Cursor for quick plan"),
        ("Minute 3", "Run PR/quality skill"),
        ("Minute 4", "Update task status safely"),
        ("Minute 5", "Save one useful prompt"),
    ]
    x = 200
    y = 350
    w = 380
    h = 540
    gap = 60
    for idx, (minute, desc) in enumerate(steps):
        box = (x, y, x + w, y + h)
        draw.rounded_rectangle(box, radius=22, fill=(236, 247, 255), outline=(125, 157, 208), width=4)
        draw.text((x + 46, y + 56), minute, font=load_font(48, bold=True), fill=(28, 51, 93))
        draw.text((x + 36, y + 178), fill(desc, width=16), font=load_font(36), fill=(41, 61, 95))
        if idx < len(steps) - 1:
            draw_arrow(draw, (x + w, y + h // 2), (x + w + gap - 16, y + h // 2), width=9)
        x += w + gap

    draw.rounded_rectangle((260, 970, 2140, 1220), radius=18, fill=(236, 252, 243), outline=(126, 180, 149), width=3)
    draw.text((320, 1040), "Small daily habit -> faster team adoption and better quality.", font=load_font(42, bold=True), fill=(29, 94, 62))

    save_image(image, path)


def generate_images() -> dict:
    ASSETS_DIR.mkdir(exist_ok=True)
    files = {
        "cover": ASSETS_DIR / "cover_visual.png",
        "top_mcp": ASSETS_DIR / "top_mcp_tools.png",
        "top_skills": ASSETS_DIR / "top_agent_skills.png",
        "architecture": ASSETS_DIR / "simple_architecture.png",
        "prompt_formula": ASSETS_DIR / "prompt_formula.png",
        "do_dont": ASSETS_DIR / "do_dont.png",
        "common_errors": ASSETS_DIR / "common_errors_fixes.png",
        "five_min_routine": ASSETS_DIR / "five_min_routine.png",
        "settings_screen": ASSETS_DIR / "cursor_settings_screen.png",
        "mcp_json_screen": ASSETS_DIR / "cursor_mcp_json_screen.png",
        "status_screen": ASSETS_DIR / "cursor_status_screen.png",
        "tutorial_path": ASSETS_DIR / "tutorial_path.png",
        "daily_workflow": ASSETS_DIR / "daily_workflow.png",
        "risk_controls": ASSETS_DIR / "risk_controls_simple.png",
        "roadmap": ASSETS_DIR / "roadmap_30_60_90.png",
    }

    create_cover_visual(files["cover"])
    create_top_mcp_visual(files["top_mcp"])
    create_top_skills_visual(files["top_skills"])
    create_architecture_visual(files["architecture"])
    create_prompt_formula_visual(files["prompt_formula"])
    create_do_dont_visual(files["do_dont"])
    create_common_errors_visual(files["common_errors"])
    create_five_min_routine_visual(files["five_min_routine"])
    create_cursor_settings_screen(files["settings_screen"])
    create_cursor_mcp_json_screen(files["mcp_json_screen"])
    create_cursor_connection_status_screen(files["status_screen"])
    create_tutorial_path_visual(files["tutorial_path"])
    create_daily_workflow_visual(files["daily_workflow"])
    create_risk_controls_visual(files["risk_controls"])
    create_roadmap_visual(files["roadmap"])
    return files


def style_title(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 40, 74)
    p.alignment = PP_ALIGN.LEFT


def add_title_block(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(12.2), Inches(0.72))
    style_title(title_box, title)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.48), Inches(0.82), Inches(12.2), Inches(0.46))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(68, 86, 120)


def add_bullets(slide, lines, x=0.55, y=1.4, w=5.4, h=5.6, level0_size=21, level1_size=17) -> None:
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(lines):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(level0_size if level == 0 else level1_size)
        p.font.color.rgb = RGBColor(35, 53, 84)
        p.space_after = Pt(8)


def add_image(slide, image_path: Path, x=6.0, y=1.25, w=7.1) -> None:
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w))


def add_code_block(slide, code_lines, x=0.55, y=1.75, w=12.2, h=5.2, font_size=14) -> None:
    block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor(22, 31, 49)
    block.line.color.rgb = RGBColor(22, 31, 49)
    tf = block.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(code_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "DejaVu Sans Mono"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(228, 236, 252)
        p.space_after = Pt(1)
        p.space_before = Pt(0)


def build_participant_presentation(images: dict, output_file: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) Cover
    slide = prs.slides.add_slide(blank)
    add_title_block(
        slide,
        "Solution Assessment: Cursor AI Capabilities and Integration Review",
        "Spoon-feed training edition: simple English, clear steps, and full tutorial.",
    )
    add_bullets(
        slide,
        [
            "What this deck gives you:",
            ("- Top MCP tools for development", 1),
            ("- Top agent skills to start fast", 1),
            ("- Step-by-step setup in Cursor", 1),
            ("- Commands your team can copy and run", 1),
            f"Updated: {datetime.now().strftime('%Y-%m-%d')}",
        ],
        x=0.6,
        y=1.55,
        w=5.2,
        h=4.8,
    )
    add_image(slide, images["cover"], x=5.95, y=1.22, w=7.15)

    # 2) Learning goals
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "What your team will learn (easy words)")
    add_bullets(
        slide,
        [
            "1) What MCP is and why it helps developers.",
            "2) Which MCP tools give the fastest value.",
            "3) How to connect Jira, Figma, and Bitbucket in Cursor.",
            "4) How to create reusable agent skills.",
            "5) How to use safe controls so automation stays trusted.",
        ],
        x=0.9,
        y=1.55,
        w=11.2,
        h=5.6,
        level0_size=26,
    )

    # 3) Agenda
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Agenda (spoon-feed order)")
    add_bullets(
        slide,
        [
            "Part A: Top MCP tools for development",
            "Part B: Top agent skills for developers",
            "Part C: Prompt formula + Do/Don't quick guide",
            "Part D: Step-by-step Cursor setup tutorial",
            "Part E: Common errors and quick fixes",
            "Part F: Daily routine + risk controls + roadmap",
        ],
        x=0.95,
        y=1.55,
        w=11.0,
        h=5.6,
        level0_size=25,
    )

    # 4) Top MCP tools
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Top MCP tools that help development teams most")
    add_bullets(
        slide,
        [
            "Pick these first in your pilot:",
            ("1) Bitbucket/Git MCP", 1),
            ("2) Jira MCP", 1),
            ("3) Figma MCP", 1),
            ("4) Docs MCP", 1),
            ("5) CI/CD MCP", 1),
            ("6) Database MCP", 1),
        ],
    )
    add_image(slide, images["top_mcp"], x=5.95, y=1.23, w=7.15)

    # 5) Top skills
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Top agent skills to create first")
    add_bullets(
        slide,
        [
            "Build a small skill library in week 1:",
            ("- jira-ticket-triage", 1),
            ("- figma-handoff", 1),
            ("- pr-quality-check", 1),
            ("- release-note-writer", 1),
            ("- bug-root-cause", 1),
            ("- test-case-generator", 1),
        ],
    )
    add_image(slide, images["top_skills"], x=5.95, y=1.23, w=7.15)

    # 6) Architecture
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "How the connection works (simple architecture)")
    add_bullets(
        slide,
        [
            "Tools -> MCP layer -> Cursor -> Governance",
            "MCP gives one common interface for each tool.",
            "Cursor uses prompts + skills to call those tools.",
            "Team rules control access, approvals, and logs.",
        ],
    )
    add_image(slide, images["architecture"], x=5.95, y=1.23, w=7.15)

    # 7) Prompt formula
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Prompt formula cheat sheet")
    add_bullets(
        slide,
        [
            "Use this every time:",
            ("Context: where and what", 1),
            ("Task: what output you need", 1),
            ("Constraints: limits and rules", 1),
            ("Output format: bullet/table/json", 1),
        ],
    )
    add_image(slide, images["prompt_formula"], x=5.95, y=1.23, w=7.15)

    # 8) Do and Don't
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Do and Don't (simple rules)")
    add_bullets(
        slide,
        [
            "Use this for new team members:",
            ("Do: read-only first, then safe write with approval", 1),
            ("Do: save good prompts as skills", 1),
            ("Don't: use admin tokens", 1),
            ("Don't: auto-merge without review", 1),
        ],
    )
    add_image(slide, images["do_dont"], x=5.95, y=1.23, w=7.15)

    # 9) Tutorial map
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Tutorial map: follow these 8 steps")
    add_bullets(
        slide,
        [
            "Trainer flow for one session (about 90 minutes):",
            ("10 min demo + 25 min pair lab + 15 min review", 1),
            ("Repeat for Jira, Figma, and Bitbucket", 1),
            ("End with one reusable skill per developer", 1),
        ],
    )
    add_image(slide, images["tutorial_path"], x=5.95, y=1.23, w=7.15)

    # 10) Step 1 commands
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 1: Prepare local machine (copy and run)")
    add_code_block(
        slide,
        [
            "mkdir -p ~/cursor-mcp-training/{servers,skills,logs} && cd ~/cursor-mcp-training",
            "python3 -m venv .venv",
            "source .venv/bin/activate",
            "python -m pip install --upgrade pip",
            "pip install mcp httpx python-dotenv pyyaml",
            "python --version && node --version && npm --version",
            "sudo apt-get update && sudo apt-get install -y jq",
            "",
            "# Create environment file for tokens",
            "cat > .env <<'EOF'",
            "JIRA_BASE_URL=https://your-company.atlassian.net",
            "JIRA_EMAIL=you@company.com",
            "JIRA_API_TOKEN=<jira_token>",
            "FIGMA_TOKEN=<figma_token>",
            "BITBUCKET_WORKSPACE=<workspace>",
            "BITBUCKET_USERNAME=<username>",
            "BITBUCKET_APP_PASSWORD=<app_password>",
            "EOF",
        ],
        y=1.55,
        h=5.8,
        font_size=14,
    )

    # 11) Step 2 screenshot
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 2: Open Cursor settings and go to MCP")
    add_bullets(
        slide,
        [
            "Click Settings in Cursor.",
            "Open Features tab.",
            "Open MCP section.",
            "Turn MCP ON.",
            "Click Open mcp.json.",
        ],
        x=0.55,
        y=1.55,
        w=4.8,
        h=5.5,
        level0_size=25,
    )
    add_image(slide, images["settings_screen"], x=5.35, y=1.22, w=7.75)

    # 12) Step 3 screenshot + config
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 3: Add mcp.json config in Cursor")
    add_bullets(
        slide,
        [
            "Paste this config in .cursor/mcp.json:",
            "Use your own server script paths.",
            "Save file and restart Cursor.",
        ],
        x=0.55,
        y=1.4,
        w=5.2,
        h=1.8,
        level0_size=22,
    )
    add_code_block(
        slide,
        [
            "{",
            '  "mcpServers": {',
            '    "jira": {"command":"python","args":["servers/jira_server.py"],"envFile":".env"},',
            '    "figma": {"command":"python","args":["servers/figma_server.py"],"envFile":".env"},',
            '    "bitbucket": {"command":"python","args":["servers/bitbucket_server.py"],"envFile":".env"}',
            "  }",
            "}",
        ],
        x=0.55,
        y=3.0,
        w=5.2,
        h=3.95,
        font_size=13,
    )
    add_image(slide, images["mcp_json_screen"], x=5.95, y=1.22, w=7.15)

    # 13) Step 4 token file
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 4: Add tokens in .env file")
    add_code_block(
        slide,
        [
            "cat > .env <<'EOF'",
            "JIRA_BASE_URL=https://your-company.atlassian.net",
            "JIRA_EMAIL=you@company.com",
            "JIRA_API_TOKEN=<jira_token>",
            "FIGMA_TOKEN=<figma_token>",
            "BITBUCKET_WORKSPACE=<workspace>",
            "BITBUCKET_USERNAME=<username>",
            "BITBUCKET_APP_PASSWORD=<app_password>",
            "EOF",
            "",
            "# Keep .env out of git",
            "echo '.env' >> .gitignore",
        ],
        y=1.9,
        h=4.95,
        font_size=15,
    )
    add_bullets(
        slide,
        [
            "Important:",
            ("Never commit .env to repository.", 1),
            ("Rotate tokens every 90 days.", 1),
        ],
        x=0.65,
        y=6.95,
        w=12.0,
        h=0.5,
        level0_size=17,
        level1_size=15,
    )

    # 14) Step 5 Jira
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 5: Test Jira connection")
    add_bullets(
        slide,
        [
            "Run test command:",
            ("curl -s -u \"$JIRA_EMAIL:$JIRA_API_TOKEN\" \"$JIRA_BASE_URL/rest/api/3/myself\" | jq '.displayName'", 1),
            "Then ask in Cursor:",
            ('"Using Jira MCP, summarize PROJ-101 and draft acceptance criteria."', 1),
        ],
        x=0.55,
        y=1.55,
        w=5.2,
        h=5.2,
        level0_size=22,
        level1_size=15,
    )
    add_image(slide, images["status_screen"], x=5.95, y=1.22, w=7.15)

    # 15) Step 6 Figma
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 6: Test Figma connection")
    add_code_block(
        slide,
        [
            "source .venv/bin/activate && source .env",
            "export FIGMA_FILE_KEY=<file_key>",
            "curl -s -H \"X-Figma-Token: $FIGMA_TOKEN\" \"https://api.figma.com/v1/files/$FIGMA_FILE_KEY\" | jq '.name'",
            "",
            "# Prompt in Cursor",
            "Using Figma MCP, extract design tokens and map components to frontend stories.",
        ],
        x=0.55,
        y=1.95,
        w=5.2,
        h=4.9,
        font_size=14,
    )
    add_image(slide, images["status_screen"], x=5.95, y=1.22, w=7.15)

    # 16) Step 7 Bitbucket
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 7: Test Bitbucket connection")
    add_code_block(
        slide,
        [
            "source .venv/bin/activate && source .env",
            "curl -s -u \"$BITBUCKET_USERNAME:$BITBUCKET_APP_PASSWORD\" \\",
            "  \"https://api.bitbucket.org/2.0/repositories/$BITBUCKET_WORKSPACE\" | jq '.values[0].full_name'",
            "",
            "# Prompt in Cursor",
            "Using Bitbucket MCP, summarize PR #123 and create review checklist.",
        ],
        x=0.55,
        y=1.95,
        w=5.2,
        h=4.9,
        font_size=14,
    )
    add_image(slide, images["status_screen"], x=5.95, y=1.22, w=7.15)

    # 17) Step 8 skill creation
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Step 8: Build your first agent skill")
    add_code_block(
        slide,
        [
            "mkdir -p skills/jira-ticket-triage && cd skills/jira-ticket-triage",
            "cat > skill.yaml <<'YAML'",
            "name: jira-ticket-triage",
            "version: 1.0.0",
            "description: Triage issue and return action plan",
            "tools: [jira.search, jira.get_issue]",
            "YAML",
            "",
            "cat > prompts.md <<'MD'",
            "# Inputs: issue_key, team_context, done_definition",
            "# Output: summary, acceptance criteria, test cases",
            "MD",
            "",
            "echo '[{\"input\":\"PROJ-101\",\"assert_contains\":[\"Summary\",\"Acceptance\"]}]' > tests.json",
        ],
        y=1.75,
        h=5.5,
        font_size=14,
    )

    # 18) Common errors
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Common errors and fixes (use in live training)")
    add_bullets(
        slide,
        [
            "Most common issues:",
            ("401: token expired or wrong", 1),
            ("403: missing project/file permission", 1),
            ("404: wrong key or URL", 1),
            ("MCP not showing: restart Cursor after save", 1),
        ],
        level0_size=20,
        level1_size=16,
    )
    add_image(slide, images["common_errors"], x=5.95, y=1.23, w=7.15)

    # 19) Five-minute routine
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "5-minute daily routine (adoption booster)")
    add_bullets(
        slide,
        [
            "Do this daily:",
            ("Minute 1: pick top task", 1),
            ("Minute 2: ask Cursor for plan", 1),
            ("Minute 3: run quality skill", 1),
            ("Minute 4: update status", 1),
            ("Minute 5: save one useful prompt", 1),
        ],
        level0_size=20,
        level1_size=16,
    )
    add_image(slide, images["five_min_routine"], x=5.95, y=1.23, w=7.15)

    # 20) Daily workflow
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Daily workflow your team can follow")
    add_bullets(
        slide,
        [
            "Keep it simple:",
            ("Morning plan -> Build -> Review -> Close", 1),
            ("Save good prompts as team skills", 1),
            ("Do weekly quality review", 1),
        ],
    )
    add_image(slide, images["daily_workflow"], x=5.95, y=1.23, w=7.15)

    # 21) Risk controls
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Risk Management and Governance Controls (simple)")
    add_bullets(
        slide,
        [
            "Four must-have controls:",
            ("1) Low-access tokens", 1),
            ("2) Human approval for write actions", 1),
            ("3) Audit logs for every automation", 1),
            ("4) Weekly KPI and failure review", 1),
        ],
    )
    add_image(slide, images["risk_controls"], x=5.95, y=1.23, w=7.15)

    # 22) Roadmap
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "30-60-90 day adoption roadmap")
    add_bullets(
        slide,
        [
            "0-30 days: connect tools and train first squad.",
            "31-60 days: publish top skills and run office hours.",
            "61-90 days: scale to more teams and report KPI impact.",
            "Target: 20% faster delivery with safe controls.",
        ],
    )
    add_image(slide, images["roadmap"], x=5.95, y=1.23, w=7.15)

    # 23) Final checklist
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Final checklist for trainer and team")
    add_bullets(
        slide,
        [
            "Before session:",
            ("Tokens ready and mcp.json template prepared", 1),
            ("Pilot repo, Jira project, and Figma file selected", 1),
            "During session:",
            ("Everyone completes one end-to-end workflow", 1),
            ("Everyone creates at least one simple skill", 1),
            "After session:",
            ("Review KPIs and improve weak prompts weekly", 1),
            "",
            "Q&A",
        ],
        x=0.9,
        y=1.55,
        w=5.2,
        h=5.8,
    )
    add_image(slide, images["cover"], x=5.95, y=1.22, w=7.15)

    prs.save(output_file)


def build_trainer_presentation(images: dict, output_file: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) Cover
    slide = prs.slides.add_slide(blank)
    add_title_block(
        slide,
        "Trainer Deck (45 min): Cursor AI + MCP Enablement",
        "Facilitator version with talking points, pacing, and troubleshooting cues.",
    )
    add_bullets(
        slide,
        [
            "Audience: developers and team leads",
            "Goal: connect tools, run safe workflows, publish first skills",
            "Format: short demo + labs + troubleshooting",
            f"Updated: {datetime.now().strftime('%Y-%m-%d')}",
        ],
        x=0.6,
        y=1.55,
        w=5.2,
        h=4.8,
    )
    add_image(slide, images["cover"], x=5.95, y=1.22, w=7.15)

    # 2) Run of show
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "45-minute trainer run-of-show")
    add_bullets(
        slide,
        [
            "0-5 min: explain MCP in simple terms",
            "5-12 min: show Top MCP and Top Skills",
            "12-22 min: live setup in Cursor settings",
            "22-34 min: Jira/Figma/Bitbucket connection tests",
            "34-40 min: build one simple agent skill",
            "40-45 min: common errors, Q&A, next steps",
        ],
        x=0.8,
        y=1.55,
        w=11.3,
        h=5.7,
        level0_size=24,
    )

    # 3) What to teach first
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "What to teach first (order matters)")
    add_bullets(
        slide,
        [
            "1) Top MCP tools",
            "2) Top agent skills",
            "3) Prompt formula",
            "4) Do/Don't safe rules",
            "5) Setup steps and command practice",
            "6) Common errors and quick fixes",
        ],
        x=0.95,
        y=1.55,
        w=11.0,
        h=5.6,
        level0_size=25,
    )

    # 4) Top MCP
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Top MCP tools for development teams")
    add_bullets(
        slide,
        [
            "Use these in your pilot first:",
            ("Bitbucket/Git, Jira, Figma", 1),
            ("Docs MCP, CI/CD MCP, Database MCP", 1),
            "Reason: fast wins and clear ROI",
        ],
    )
    add_image(slide, images["top_mcp"], x=5.95, y=1.23, w=7.15)

    # 5) Top Skills
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Top starter agent skills")
    add_bullets(
        slide,
        [
            "Required first set:",
            ("jira-ticket-triage", 1),
            ("figma-handoff", 1),
            ("pr-quality-check", 1),
            ("release-note-writer", 1),
            ("bug-root-cause", 1),
            ("test-case-generator", 1),
        ],
    )
    add_image(slide, images["top_skills"], x=5.95, y=1.23, w=7.15)

    # 6) Prompt formula
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Prompt formula (repeat every demo)")
    add_bullets(
        slide,
        [
            "Trainer talking line:",
            ('"Context + Task + Constraints + Output format"', 1),
            "Ask participants to use this format in all labs.",
            "This reduces confusion and bad outputs.",
        ],
    )
    add_image(slide, images["prompt_formula"], x=5.95, y=1.23, w=7.15)

    # 7) Do/Don't
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Do and Don't (safety first)")
    add_bullets(
        slide,
        [
            "Call this out clearly before labs:",
            ("Do: use read-only first", 1),
            ("Do: request approval before write actions", 1),
            ("Don't: use admin tokens", 1),
            ("Don't: skip review/checklist", 1),
        ],
    )
    add_image(slide, images["do_dont"], x=5.95, y=1.23, w=7.15)

    # 8) Setup commands
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Lab setup commands (copy and run)")
    add_code_block(
        slide,
        [
            "mkdir -p ~/cursor-mcp-training/{servers,skills,logs} && cd ~/cursor-mcp-training",
            "python3 -m venv .venv && source .venv/bin/activate",
            "python -m pip install --upgrade pip",
            "pip install mcp httpx python-dotenv pyyaml",
            "",
            "cat > .env <<'EOF'",
            "JIRA_BASE_URL=https://your-company.atlassian.net",
            "JIRA_EMAIL=you@company.com",
            "JIRA_API_TOKEN=<jira_token>",
            "FIGMA_TOKEN=<figma_token>",
            "BITBUCKET_WORKSPACE=<workspace>",
            "BITBUCKET_USERNAME=<username>",
            "BITBUCKET_APP_PASSWORD=<app_password>",
            "EOF",
        ],
        y=1.75,
        h=5.45,
        font_size=14,
    )

    # 9) Cursor settings screenshot
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Live demo: Cursor settings -> MCP")
    add_bullets(
        slide,
        [
            "Step 1: open Settings",
            "Step 2: Features -> MCP",
            "Step 3: switch MCP to ON",
            "Step 4: open mcp.json",
        ],
        x=0.55,
        y=1.55,
        w=4.7,
        h=5.4,
        level0_size=24,
    )
    add_image(slide, images["settings_screen"], x=5.35, y=1.22, w=7.75)

    # 10) mcp.json screenshot
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Live demo: mcp.json content")
    add_code_block(
        slide,
        [
            "{",
            '  "mcpServers": {',
            '    "jira": {"command":"python","args":["servers/jira_server.py"],"envFile":".env"},',
            '    "figma": {"command":"python","args":["servers/figma_server.py"],"envFile":".env"},',
            '    "bitbucket": {"command":"python","args":["servers/bitbucket_server.py"],"envFile":".env"}',
            "  }",
            "}",
            "",
            "# Save and restart Cursor",
        ],
        x=0.55,
        y=2.1,
        w=5.2,
        h=4.8,
        font_size=13,
    )
    add_image(slide, images["mcp_json_screen"], x=5.95, y=1.22, w=7.15)

    # 11) Connection tests
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Live demo: test connections (Jira/Figma/Bitbucket)")
    add_code_block(
        slide,
        [
            "# Jira",
            "curl -s -u \"$JIRA_EMAIL:$JIRA_API_TOKEN\" \"$JIRA_BASE_URL/rest/api/3/myself\" | jq '.displayName'",
            "",
            "# Figma",
            "curl -s -H \"X-Figma-Token: $FIGMA_TOKEN\" \"https://api.figma.com/v1/files/$FIGMA_FILE_KEY\" | jq '.name'",
            "",
            "# Bitbucket",
            "curl -s -u \"$BITBUCKET_USERNAME:$BITBUCKET_APP_PASSWORD\" \"https://api.bitbucket.org/2.0/repositories/$BITBUCKET_WORKSPACE\" | jq '.values[0].full_name'",
        ],
        x=0.55,
        y=2.0,
        w=5.2,
        h=4.9,
        font_size=12,
    )
    add_image(slide, images["status_screen"], x=5.95, y=1.22, w=7.15)

    # 12) Common errors
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Live support slide: common errors and fixes")
    add_bullets(
        slide,
        [
            "Use this during Q&A:",
            ("401 -> token issue", 1),
            ("403 -> permission issue", 1),
            ("404 -> wrong key/url", 1),
            ("MCP not listed -> restart Cursor", 1),
        ],
    )
    add_image(slide, images["common_errors"], x=5.95, y=1.23, w=7.15)

    # 13) 5-minute routine
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Coach this habit: 5-minute daily routine")
    add_bullets(
        slide,
        [
            "Ask every participant to follow this daily.",
            "This is the easiest way to keep adoption active.",
            "Review weekly and celebrate improvements.",
        ],
    )
    add_image(slide, images["five_min_routine"], x=5.95, y=1.23, w=7.15)

    # 14) Risk controls and roadmap
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Risk controls and 30-60-90 rollout")
    add_bullets(
        slide,
        [
            "Must-have controls:",
            ("Low-access tokens + human approvals", 1),
            ("Audit logs + weekly KPI review", 1),
            "Rollout: 0-30 connect, 31-60 standardize, 61-90 scale.",
        ],
        level0_size=19,
        level1_size=16,
    )
    add_image(slide, images["roadmap"], x=5.95, y=1.23, w=7.15)

    # 15) Final trainer checklist
    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "Trainer checklist (before, during, after)")
    add_bullets(
        slide,
        [
            "Before: test tokens and demo environment.",
            "During: keep prompts simple and repeat formula.",
            "During: pause at each error and show quick fix.",
            "After: share participant deck and runbook.",
            "After: schedule weekly office hours.",
            "",
            "Q&A",
        ],
        x=0.8,
        y=1.55,
        w=5.5,
        h=5.7,
    )
    add_image(slide, images["cover"], x=5.95, y=1.22, w=7.15)

    prs.save(output_file)


def main() -> None:
    images = generate_images()
    build_participant_presentation(images, PARTICIPANT_OUTPUT_FILE)
    copyfile(PARTICIPANT_OUTPUT_FILE, LEGACY_OUTPUT_FILE)
    build_trainer_presentation(images, TRAINER_OUTPUT_FILE)
    print(f"Created participant deck: {PARTICIPANT_OUTPUT_FILE}")
    print(f"Created trainer deck: {TRAINER_OUTPUT_FILE}")
    print(f"Updated legacy deck: {LEGACY_OUTPUT_FILE}")
    print(f"Assets directory: {ASSETS_DIR}")


if __name__ == "__main__":
    main()
